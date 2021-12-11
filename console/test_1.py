from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR



path_2 = 'g:\\temp\\test2.pptx'
path_3 = 'g:\\temp\\test3.pptx'

prs = Presentation(path_2)
slide = prs.slides[0]
for i, shape in enumerate(slide.shapes):
    print(i, shape.text)

left = top = Inches(1)
width = Inches(2)
height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
txBox.alignment = PP_ALIGN.CENTER
tf = txBox.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Hello World"
font = run.font
font.name = 'Tahoma'
font.size = Pt(12)
font.bold = True

first = slide.shapes[0]
first.left = Inches(12.33)
first.width = Inches(1)
first.height = Inches(1)
first.top = Inches(6.5)


prs.save(path_3)