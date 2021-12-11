from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class MyPresentation:
    WIDTH = 13.33
    HEIGHT = 7.5

    def __init__(self, path_in, path_out):
        self.path_out = path_out
        self.p = Presentation(path_in)
        self.slide = self.p.slides[0]
        self.shapes = self.slide.shapes
        # map user-given name with shape
        self.shape = dict()

    def map(self, li):
        assert type(li) in {tuple, list, set}
        assert len(li) == len(self.shapes)
        for i, name in enumerate(li):
            self.shape[name] = self.shapes[i]

    def update(self, name, parent, ratios):
        left, top, width, height = parent.get_coor(ratios)
        self.shape[name].left = Inches(self.WIDTH / 100 * left)
        self.shape[name].top = Inches(self.HEIGHT / 100 * top)
        self.shape[name].width = Inches(self.WIDTH / 100 * width)
        self.shape[name].height = Inches(self.HEIGHT / 100 * height)

    def add_title(self, name, text, font_size, coor):

        textbox = self.slide.shapes.add_textbox(left, top, width, height)
        self.shape[name] = self.shapes[-1]
        tf = textbox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        parag = tf.paragraphs[0]
        parag.alignment = PP_ALIGN.CENTER
        run = parag.add_run()
        run.text = text
        font = run.font
        font.name = 'Tahoma'
        font.size = Pt(font_size)
        font.bold = True

    def save(self):
        self.p.save(self.path_out)

    @classmethod
    def to_inches(cls, coor):
        left = Inches(cls.WIDTH / 100 * coor[0])
        top = Inches(cls.HEIGHT / 100 * coor[1])
        width = Inches(cls.WIDTH / 100 * coor[2])
        height = Inches(cls.HEIGHT / 100 * coor[3])


from c_group_ratio import GroupRatio

path_in = 'g:\\temp\\test4.pptx'
path_out = 'g:\\temp\\test6.pptx'

p = MyPresentation(path_in, path_out)
p.map(list('abcd'))
g_slide = GroupRatio()
g_top = GroupRatio(parent=g_slide, ratios=(0, 0, 100, 40))
g_bottom = GroupRatio(parent=g_slide, ratios=(0, 40, 100, 60))
p.update('a', g_top, ratios=(10, 10, 20, 89))
p.update('b', g_top, ratios=(40, 10, 40, 100))
p.update('c', g_bottom, ratios=(10, 10, 50, 80))
p.update('d', g_bottom, ratios=(70, 10, 20 ,80))
p.add_title(name='e', text='My Title', font_size=20,
            coor=g_top.get_coor(ratios=(10, 10, 80, 20)))
p.save()